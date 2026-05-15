from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0D, 0x1B, 0x2A)   # dark bg
BLUE    = RGBColor(0x1A, 0x73, 0xE8)   # accent
CYAN    = RGBColor(0x00, 0xC2, 0xFF)   # highlight
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xE8, 0xEC, 0xF0)
DGRAY   = RGBColor(0x60, 0x70, 0x80)
GREEN   = RGBColor(0x1E, 0xC8, 0x7E)
ORANGE  = RGBColor(0xFF, 0x8C, 0x00)
RED     = RGBColor(0xFF, 0x4C, 0x4C)
YELLOW  = RGBColor(0xFF, 0xD7, 0x00)
CARD_BG = RGBColor(0x16, 0x2A, 0x42)   # card background


# ── Helpers ───────────────────────────────────────────────────────────────────
def blank_layout(prs):
    return prs.slide_layouts[6]   # truly blank

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def full_bg(slide, color=NAVY):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=color)

def slide_header(slide, title, subtitle=None):
    # top accent bar
    add_rect(slide, 0, 0, 13.33, 0.08, fill=BLUE)
    add_text(slide, title, 0.5, 0.18, 12, 0.65, size=30, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.5, 0.82, 12, 0.42, size=15, color=CYAN, italic=True)

def card(slide, l, t, w, h, fill=CARD_BG, radius=False):
    return add_rect(slide, l, t, w, h, fill=fill)

def icon_card(slide, l, t, icon, title, body, w=2.8, h=1.55):
    card(slide, l, t, w, h)
    add_text(slide, icon,  l+0.12, t+0.1,  0.55, 0.55, size=22, align=PP_ALIGN.CENTER)
    add_text(slide, title, l+0.7,  t+0.12, w-0.8, 0.38, size=13, bold=True, color=CYAN)
    add_text(slide, body,  l+0.7,  t+0.48, w-0.8, 0.95, size=11, color=LGRAY)

def bullet_box(slide, l, t, w, h, items, icon_color=CYAN, fill=CARD_BG, title=None):
    card(slide, l, t, w, h, fill=fill)
    y = t + 0.15
    if title:
        add_text(slide, title, l+0.2, y, w-0.4, 0.35, size=13, bold=True, color=CYAN)
        y += 0.38
    for item in items:
        add_rect(slide, l+0.18, y+0.09, 0.06, 0.06, fill=icon_color)
        add_text(slide, item, l+0.35, y, w-0.55, 0.32, size=11, color=WHITE)
        y += 0.32

def table_slide_rows(slide, l, t, w, col_ws, headers, rows, hdr_fill=BLUE):
    x = l
    row_h = 0.34
    for i, (hdr, cw) in enumerate(zip(headers, col_ws)):
        add_rect(slide, x, t, cw, row_h, fill=hdr_fill)
        add_text(slide, hdr, x+0.08, t+0.04, cw-0.1, row_h-0.06,
                 size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += cw
    for ri, row in enumerate(rows):
        x = l
        fill = CARD_BG if ri % 2 == 0 else RGBColor(0x1A, 0x2F, 0x48)
        for ci, (cell, cw) in enumerate(zip(row, col_ws)):
            add_rect(slide, x, t + row_h*(ri+1), cw, row_h, fill=fill)
            col = GREEN if str(cell).startswith("✓") else (RED if str(cell).startswith("✗") else WHITE)
            add_text(slide, str(cell), x+0.08, t+row_h*(ri+1)+0.04,
                     cw-0.1, row_h-0.06, size=10, color=col, align=PP_ALIGN.CENTER)
            x += cw


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
full_bg(sl)
# large gradient-feel accent band
add_rect(sl, 0, 2.6, 13.33, 2.8, fill=RGBColor(0x10, 0x22, 0x36))
add_rect(sl, 0, 2.6, 0.12, 2.8, fill=BLUE)

add_text(sl, "Order-Based Planning System", 0.5, 2.8, 12.3, 1.0,
         size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Intelligent Manufacturing Scheduling using Genetic Algorithm Optimization",
         0.5, 3.75, 12.3, 0.6, size=18, color=CYAN, align=PP_ALIGN.CENTER, italic=True)
add_text(sl, "PrecisionParts Ltd. — Automotive Components Manufacturer",
         0.5, 4.45, 12.3, 0.45, size=14, color=DGRAY, align=PP_ALIGN.CENTER)

# decorative dots
for xi, yi in [(1.2,1.1),(2.5,0.6),(4.0,1.4),(10.5,1.0),(11.8,0.55),(12.3,1.6)]:
    add_rect(sl, xi, yi, 0.09, 0.09, fill=BLUE)

add_text(sl, "Slide 1 of 11", 0.3, 7.1, 3, 0.3, size=9, color=DGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE BUSINESS PROBLEM
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
full_bg(sl)
slide_header(sl, "The Business Problem",
             "What challenge does Order-Based Planning solve?")

# central question box
add_rect(sl, 1.0, 1.55, 11.33, 1.1, fill=RGBColor(0x0A, 0x2A, 0x50))
add_rect(sl, 1.0, 1.55, 0.1, 1.1, fill=CYAN)
add_text(sl, '❝  Given limited machines and materials, which orders do we produce this week — '
             'and what do we delay to minimize total penalty cost?  ❞',
         1.2, 1.65, 10.9, 0.9, size=14, color=WHITE, italic=True, align=PP_ALIGN.CENTER)

pain_points = [
    ("⚠", "Capacity Overload",   "Multiple orders compete for the same CNC machine or assembly bay in the same week"),
    ("📦", "Material Shortages",  "Steel rods or rubber seals run out mid-week; some orders can't be fulfilled on time"),
    ("📅", "Promise Date Conflicts", "10 orders all promise delivery on May 16 — physically impossible to honour all"),
    ("💸", "Penalty Blind Spots",  "Delaying a VIP customer (Toyota) costs 10× more than delaying a low-priority garage"),
]
for i, (ico, ttl, desc) in enumerate(pain_points):
    col = i % 2
    row = i // 2
    icon_card(sl, 0.5 + col*6.42, 2.88 + row*1.7, ico, ttl, desc, w=6.15, h=1.58)

add_text(sl, "→ Manual spreadsheet planning misses interactions between orders, capacity and components.",
         0.5, 6.35, 12.3, 0.45, size=12, color=ORANGE, italic=True, align=PP_ALIGN.CENTER)
add_text(sl, "Slide 2 of 11", 0.3, 7.1, 3, 0.3, size=9, color=DGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — MEET PRECISIONPARTS LTD.
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
full_bg(sl)
slide_header(sl, "Meet PrecisionParts Ltd.",
             "Our example: an automotive components manufacturer")

add_text(sl, "🏭", 0.5, 1.4, 1.0, 0.8, size=40)
add_text(sl, "PrecisionParts Ltd. manufactures Engine Blocks, Brake Assemblies, and Exhaust Pipes\n"
             "for OEM customers. They receive ~30 orders per month across 3 factory locations.",
         1.5, 1.45, 11.0, 0.75, size=13, color=LGRAY)

# 3 columns
cols = [
    ("🔧 Products", ["Engine Block V6 — $1,200", "Brake Assembly — $340", "Exhaust Pipe — $180"],
     "What they manufacture"),
    ("🏢 Customers", ["Toyota → VIP priority", "Ford → High priority", "Local Garages → Low priority"],
     "Who they sell to"),
    ("📍 Locations", ["Plant A — Detroit (USA)", "Plant B — Stuttgart (DE)", "Plant C — Chennai (IN)"],
     "Where they produce"),
]
for i, (ttl, items, sub) in enumerate(cols):
    x = 0.5 + i * 4.28
    card(sl, x, 2.5, 3.95, 3.55)
    add_text(sl, ttl, x+0.2, 2.62, 3.55, 0.44, size=14, bold=True, color=CYAN)
    add_text(sl, sub, x+0.2, 3.02, 3.55, 0.32, size=10, color=DGRAY, italic=True)
    add_rect(sl, x+0.2, 3.36, 3.55, 0.03, fill=BLUE)
    for j, item in enumerate(items):
        add_rect(sl, x+0.25, 3.52+j*0.55+0.13, 0.07, 0.07, fill=BLUE)
        add_text(sl, item, x+0.42, 3.48+j*0.55, 3.3, 0.4, size=11, color=WHITE)

add_text(sl, "Slide 3 of 11", 0.3, 7.1, 3, 0.3, size=9, color=DGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — THE 6 BUILDING BLOCKS (data model)
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
full_bg(sl)
slide_header(sl, "The 6 Building Blocks",
             "Master data that drives the planning engine")

blocks = [
    ("📦", "Products",       "What is being manufactured\n(code, name, price, lead time)"),
    ("🏢", "Customers",      "Who ordered it\n(priority: VIP / High / Low)"),
    ("⚙",  "Restrictions",   "Production resources with weekly capacity\n(CNC Line, Assembly Bay, Paint Shop)"),
    ("🔩", "Components",     "Raw materials with weekly availability\n(Steel Rods, Rubber Seals, Bolts)"),
    ("📋", "Sales Orders",   "Customer orders with promise dates,\nlinked to restrictions + components"),
    ("💸", "Penalty Rules",  "Cost per day of delay per customer priority\n(VIP: $500/day, High: $200/day)"),
]
for i, (ico, ttl, desc) in enumerate(blocks):
    col = i % 3
    row = i // 3
    x = 0.4 + col * 4.28
    y = 1.45 + row * 2.55
    card(sl, x, y, 4.1, 2.38)
    add_text(sl, ico,  x+0.15, y+0.18, 0.65, 0.65, size=26)
    add_text(sl, ttl,  x+0.8,  y+0.20, 3.0,  0.44, size=14, bold=True, color=CYAN)
    add_text(sl, desc, x+0.15, y+0.82, 3.8,  1.3,  size=11, color=LGRAY)

add_text(sl, "Slide 4 of 11", 0.3, 7.1, 3, 0.3, size=9, color=DGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — SALES ORDER ANATOMY
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
full_bg(sl)
slide_header(sl, "Anatomy of a Sales Order",
             "How a single order connects to resources and materials")

# Central order card
card(sl, 4.4, 1.5, 4.53, 2.6, fill=RGBColor(0x0A, 0x2A, 0x50))
add_rect(sl, 4.4, 1.5, 4.53, 0.08, fill=BLUE)
add_text(sl, "📋 SO-0001 — Toyota",     4.55, 1.6,  4.2, 0.42, size=14, bold=True, color=WHITE)
add_text(sl, "Product:  Engine Block V6",4.55, 2.02, 4.2, 0.32, size=11, color=LGRAY)
add_text(sl, "Quantity: 50 units",       4.55, 2.32, 4.2, 0.32, size=11, color=LGRAY)
add_text(sl, "Promise:  16 May 2025",    4.55, 2.62, 4.2, 0.32, size=11, color=LGRAY)
add_text(sl, "Priority: VIP",            4.55, 2.92, 4.2, 0.32, size=11, color=CYAN)
add_text(sl, "Penalty:  $500 / day late",4.55, 3.22, 4.2, 0.32, size=11, color=RED)

# Restrictions side
card(sl, 0.4, 1.5, 3.6, 2.6)
add_text(sl, "⚙ Restrictions Used",   0.55, 1.6,  3.3, 0.38, size=13, bold=True, color=CYAN)
for j, (rname, usage) in enumerate([
        ("CNC Machine Line A", "3 hrs / unit  →  150 hrs total"),
        ("Assembly Bay",       "1.5 hrs / unit  →  75 hrs total"),
]):
    add_rect(sl, 0.55, 2.1+j*0.65+0.12, 0.07, 0.07, fill=ORANGE)
    add_text(sl, rname,  0.72, 2.08+j*0.65, 3.1, 0.3,  size=11, bold=True, color=WHITE)
    add_text(sl, usage,  0.72, 2.36+j*0.65, 3.1, 0.28, size=10, color=LGRAY)

# Components side
card(sl, 9.33, 1.5, 3.6, 2.6)
add_text(sl, "🔩 Components Used",   9.48, 1.6,  3.3, 0.38, size=13, bold=True, color=CYAN)
for j, (cname, qty) in enumerate([
        ("Steel Rods",    "2 per unit  →  100 rods total"),
        ("Aluminum Sheet","1 per unit  →  50 sheets total"),
        ("Fastener Kit",  "4 per unit  →  200 kits total"),
]):
    add_rect(sl, 9.48, 2.1+j*0.55+0.1, 0.07, 0.07, fill=GREEN)
    add_text(sl, cname, 9.65, 2.08+j*0.55, 3.0, 0.28, size=11, bold=True, color=WHITE)
    add_text(sl, qty,   9.65, 2.34+j*0.55, 3.0, 0.26, size=10, color=LGRAY)

# Arrows
def arrow_line(sl, x1, y1, x2, y2):
    from pptx.util import Pt
    connector = sl.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = BLUE
    connector.line.width = Pt(1.5)

arrow_line(sl, 4.0, 2.8, 4.4, 2.8)
arrow_line(sl, 8.93, 2.8, 9.33, 2.8)

add_text(sl, "Each order explicitly declares which machines it occupies and which materials it consumes.\n"
             "The optimizer uses this to calculate weekly load across all orders simultaneously.",
         0.5, 4.35, 12.3, 0.75, size=12, color=LGRAY, align=PP_ALIGN.CENTER, italic=True)

add_text(sl, "Slide 5 of 11", 0.3, 7.1, 3, 0.3, size=9, color=DGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — HOW THE GENETIC ALGORITHM WORKS
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
full_bg(sl)
slide_header(sl, "The Optimization Engine",
             "Genetic Algorithm — finds the lowest-penalty schedule")

# Key question
add_rect(sl, 0.5, 1.42, 12.33, 0.72, fill=RGBColor(0x0A, 0x2A, 0x50))
add_rect(sl, 0.5, 1.42, 0.1, 0.72, fill=CYAN)
add_text(sl, "For each sales order: should we process it ON its promise date, "
             "or push it out 1, 2 … 8 weeks to relieve congestion?",
         0.72, 1.5, 12.0, 0.55, size=13, color=WHITE, italic=True)

# Steps
steps = [
    ("1", "Encode",      "Each candidate schedule = a 'chromosome'\n{SO-0001: delay 0wk, SO-0002: delay 1wk, …}"),
    ("2", "Evaluate",    "Fitness = total penalty cost\n(late delivery + capacity overrun + shortage)"),
    ("3", "Select",      "Lower-cost schedules survive\nTournament selection across 50 candidates"),
    ("4", "Crossover",   "Combine two good schedules\nto produce a better child schedule"),
    ("5", "Mutate",      "Randomly shift one order's week\nto escape local optima"),
    ("6", "Repeat",      "Evolve for 100 generations\nReturn the best schedule found"),
]
for i, (num, ttl, desc) in enumerate(steps):
    col = i % 3
    row = i // 2 if i < 6 else 1
    x = 0.5 + col * 4.28
    y = 2.35 + (i // 3) * 2.08
    card(sl, x, y, 4.1, 1.88)
    # number bubble
    add_rect(sl, x+0.18, y+0.22, 0.5, 0.5, fill=BLUE)
    add_text(sl, num,  x+0.18, y+0.22, 0.5, 0.5, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, ttl,  x+0.82, y+0.22, 3.1, 0.42, size=14, bold=True, color=CYAN)
    add_text(sl, desc, x+0.18, y+0.78, 3.75, 1.0, size=10, color=LGRAY)

add_text(sl, "Slide 6 of 11", 0.3, 7.1, 3, 0.3, size=9, color=DGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — PENALTY COST BREAKDOWN
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
full_bg(sl)
slide_header(sl, "How Penalty Cost is Calculated",
             "The fitness function the algorithm minimises")

penalty_types = [
    ("💸", "Late Delivery",    RED,    "Per order delayed past promise date",
     ["Delay days × penalty rate per day", "VIP customer: $500/day late",
      "High priority: $200/day", "Low priority: $50/day",
      "Example: Toyota SO-0001 delayed 7 days = $3,500"]),
    ("⚙",  "Capacity Overrun", ORANGE, "When a machine is over-booked in a week",
     ["Over-usage × penalty_cost_per_unit", "Soft: capacity > 0 → linear penalty",
      "Hard: capacity = 0 → ×1,000,000,000 penalty",
      "Order marked INFEASIBLE",
      "Example: CNC +40 hrs × $500 = $20,000"]),
    ("🔩", "Component Shortage", YELLOW, "When stock is insufficient for a week",
     ["Shortage units × unit_cost × 3",
      "Soft: some stock → 3× cost multiplier",
      "Hard: zero stock → ×1,000,000,000 penalty",
      "Order marked INFEASIBLE",
      "Example: 30 rods short × $45 × 3 = $4,050"]),
]
for i, (ico, ttl, col, sub, items) in enumerate(penalty_types):
    x = 0.4 + i * 4.31
    card(sl, x, 1.55, 4.15, 5.55)
    add_rect(sl, x, 1.55, 4.15, 0.08, fill=col)
    add_text(sl, ico,  x+0.18, 1.65, 0.65, 0.65, size=28)
    add_text(sl, ttl,  x+0.82, 1.7,  3.1,  0.44, size=14, bold=True, color=col)
    add_text(sl, sub,  x+0.18, 2.25, 3.8,  0.32, size=10, color=DGRAY, italic=True)
    add_rect(sl, x+0.18, 2.6, 3.8, 0.03, fill=col)
    for j, item in enumerate(items):
        add_rect(sl, x+0.22, 2.75+j*0.56+0.14, 0.06, 0.06, fill=col)
        add_text(sl, item, x+0.38, 2.73+j*0.56, 3.6, 0.4, size=10, color=WHITE)

add_text(sl, "Slide 7 of 11", 0.3, 7.1, 3, 0.3, size=9, color=DGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — OPTIMIZATION RESULTS DASHBOARD
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
full_bg(sl)
slide_header(sl, "Optimization Results Dashboard",
             "What the planner sees after running the algorithm")

# KPI cards
kpis = [
    ("70%",   "On-Time Delivery",   GREEN),
    ("7 / 10","Orders On Time",     CYAN),
    ("4.2 d", "Avg Delay",          YELLOW),
    ("$18,400","Total Penalty",     RED),
    ("2",     "Capacity Violations",ORANGE),
    ("1",     "Component Shortages",ORANGE),
]
for i, (val, lbl, col) in enumerate(kpis):
    x = 0.4 + i * 2.09
    card(sl, x, 1.45, 1.95, 1.32)
    add_rect(sl, x, 1.45, 1.95, 0.07, fill=col)
    add_text(sl, val, x, 1.6,  1.95, 0.62, size=24, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(sl, lbl, x, 2.2,  1.95, 0.48, size=9,  color=LGRAY, align=PP_ALIGN.CENTER)

# 4 result tabs description
tabs = [
    ("📋 Order Results",   "Per-order table: original vs optimised date, delay days, penalty, feasibility.\nClick any row to see components & restrictions used."),
    ("⚙ Capacity Analysis","Heat map of weekly machine utilisation.\nRed = over capacity. Drill into each restriction."),
    ("🔩 Component Analysis","Weekly stock vs demand chart.\nShortage weeks highlighted. Sorted by demand ratio."),
    ("🚨 Constraint Summary","All violations ranked by cost.\nCapacity + component + late orders in one view."),
]
for i, (ttl, desc) in enumerate(tabs):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.47
    y = 3.05 + row * 1.78
    card(sl, x, y, 6.2, 1.62)
    add_text(sl, ttl,  x+0.2, y+0.15, 5.8, 0.42, size=13, bold=True, color=CYAN)
    add_text(sl, desc, x+0.2, y+0.55, 5.8, 0.95, size=11, color=LGRAY)

add_text(sl, "Slide 8 of 11", 0.3, 7.1, 3, 0.3, size=9, color=DGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — ORDER RESULTS TABLE
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
full_bg(sl)
slide_header(sl, "Order Results — Sample Output",
             "Each row shows what changed for a sales order; click for full detail")

col_ws  = [1.55, 1.55, 1.8, 0.9, 1.55, 1.55, 0.85, 1.0, 1.0]
headers = ["Order #","Customer","Product","Priority","Original","Optimised","Delay","Penalty","Status"]
rows = [
    ["SO-0001","Toyota",    "Engine Block V6","VIP",  "16 May","16 May","✓ 0d", "$0",     "✓ On Time"],
    ["SO-0002","Ford",      "Brake Assembly", "High", "16 May","23 May","+7d",  "$1,400", "Delayed"],
    ["SO-0003","Local Garage","Exhaust Pipe", "Low",  "23 May","23 May","✓ 0d", "$0",     "✓ On Time"],
    ["SO-0004","Toyota",    "Brake Assembly", "VIP",  "23 May","30 May","+7d",  "$3,500", "Delayed"],
    ["SO-0005","Ford",      "Exhaust Pipe",   "High", "30 May","30 May","✓ 0d", "$0",     "✓ On Time"],
]
table_slide_rows(sl, 0.35, 1.52, 13.0, col_ws, headers, rows)

add_text(sl, "🖱  Click any row → modal shows Components Used, Restrictions Used, and Date of Processing",
         0.35, 4.35, 12.6, 0.44, size=12, color=CYAN, italic=True)

# mini modal mockup
card(sl, 3.5, 4.9, 6.33, 2.2, fill=RGBColor(0x10, 0x22, 0x36))
add_rect(sl, 3.5, 4.9, 6.33, 0.08, fill=BLUE)
add_text(sl, "Order SO-0001  ×", 3.65, 4.95, 5.9, 0.38, size=12, bold=True, color=WHITE)
add_text(sl, "Customer: Toyota  |  Product: Engine Block V6  |  Priority: VIP  |  Qty: 50",
         3.65, 5.35, 6.0, 0.32, size=10, color=LGRAY)
add_text(sl, "Date of Processing: 16 May 2025",
         3.65, 5.65, 6.0, 0.3, size=10, color=CYAN)
add_text(sl, "⚙ Restrictions: CNC Line A (3 hrs/unit),  Assembly Bay (1.5 hrs/unit)",
         3.65, 5.95, 6.0, 0.3, size=10, color=ORANGE)
add_text(sl, "🔩 Components:  Steel Rods (2/unit),  Aluminum Sheet (1/unit),  Fastener Kit (4/unit)",
         3.65, 6.25, 6.0, 0.3, size=10, color=GREEN)

add_text(sl, "Slide 9 of 11", 0.3, 7.1, 3, 0.3, size=9, color=DGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CAPACITY & COMPONENT ANALYSIS
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
full_bg(sl)
slide_header(sl, "Capacity & Component Analysis",
             "Visual drill-down into machine load and material stock")

# Left — capacity heat map
card(sl, 0.4, 1.52, 5.95, 5.55)
add_text(sl, "⚙ Capacity Heat Map — CNC Machine Line A", 0.58, 1.62, 5.6, 0.42, size=12, bold=True, color=CYAN)
weeks = [("W18","48%",GREEN), ("W19","72%",YELLOW), ("W20","132%",RED),
         ("W21","85%",ORANGE), ("W22","55%",GREEN), ("W23","41%",GREEN)]
for i, (wk, pct, col) in enumerate(weeks):
    x = 0.58 + i * 0.9
    add_rect(sl, x, 2.2, 0.78, 1.0, fill=col)
    add_text(sl, pct, x, 2.2,  0.78, 0.62, size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(sl, wk,  x, 2.8,  0.78, 0.32, size=9,  color=NAVY, align=PP_ALIGN.CENTER)

add_text(sl, "■ <60% OK   ■ 60-80% Moderate   ■ 80-100% High   ■ >100% Violation",
         0.58, 3.35, 5.6, 0.32, size=9, color=LGRAY)

# Violation detail
add_text(sl, "⚠ Week 20 Violation Detail", 0.58, 3.78, 5.6, 0.35, size=12, bold=True, color=RED)
add_text(sl, "Capacity Available:  200 hrs", 0.75, 4.15, 5.3, 0.3, size=11, color=WHITE)
add_text(sl, "Required:              264 hrs", 0.75, 4.45, 5.3, 0.3, size=11, color=RED)
add_text(sl, "Over-capacity:      +64 hrs", 0.75, 4.75, 5.3, 0.3, size=11, color=ORANGE)
add_text(sl, "Violation Cost:       $32,000", 0.75, 5.05, 5.3, 0.3, size=11, color=RED)
add_text(sl, "→ Solution: optimizer delayed SO-0002 by 1 week", 0.58, 5.5, 5.6, 0.42,
         size=10, color=GREEN, italic=True)

# Right — component analysis
card(sl, 6.85, 1.52, 6.08, 5.55)
add_text(sl, "🔩 Component Analysis — Steel Rods",   7.02, 1.62, 5.7, 0.42, size=12, bold=True, color=CYAN)
col_ws2 = [0.72, 1.05, 1.05, 0.88, 1.1]
hdrs2   = ["Week","Available","Required","Shortage","Status"]
rows2   = [
    ["W18","500","320","✓ 0",   "✓ OK"],
    ["W19","400","380","✓ 0",   "✓ OK"],
    ["W20","200","340","✗ 140", "Shortage"],
    ["W21","600","280","✓ 0",   "✓ OK"],
    ["W22","350","210","✓ 0",   "✓ OK"],
]
table_slide_rows(sl, 6.88, 2.15, 5.8, col_ws2, hdrs2, rows2)
add_text(sl, "Week 20 shortage (140 rods) caused orders using Steel Rods\nto shift to Week 21 where 600 units are available.",
         7.02, 5.2, 5.7, 0.75, size=11, color=LGRAY, italic=True)

add_text(sl, "Slide 10 of 11", 0.3, 7.1, 3, 0.3, size=9, color=DGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — BUSINESS VALUE & SUMMARY
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout(prs))
full_bg(sl)
slide_header(sl, "Business Value & Summary",
             "From guesswork to an optimised, defensible production schedule")

benefits = [
    ("⚡", "Speed",         "Schedule 30 orders in seconds\nwhere manual planning takes hours"),
    ("🎯", "Accuracy",      "Simultaneously respects machines,\nmaterials AND customer priorities"),
    ("💡", "Transparency",  "Every delay explained with\nexact penalty cost breakdown"),
    ("📤", "Actionability", "Excel export ready for\nshop floor and ERP upload"),
]
for i, (ico, ttl, desc) in enumerate(benefits):
    x = 0.4 + i * 3.14
    card(sl, x, 1.52, 3.0, 2.05)
    add_text(sl, ico,  x+0.15, 1.65, 0.65, 0.65, size=26)
    add_text(sl, ttl,  x+0.78, 1.68, 2.05, 0.42, size=14, bold=True, color=CYAN)
    add_text(sl, desc, x+0.15, 2.22, 2.78, 0.88, size=11, color=LGRAY)

# Flow summary
add_text(sl, "How it works — end to end", 0.4, 3.82, 12.5, 0.38, size=13, bold=True, color=WHITE)
flow = ["Enter\nOrders","Set Machine\nCapacities","Define Material\nAvailability",
        "Run\nOptimization","Review\nResults","Export\nSchedule"]
flow_cols = [GREEN, BLUE, CYAN, ORANGE, YELLOW, GREEN]
for i, (step, col) in enumerate(zip(flow, flow_cols)):
    x = 0.4 + i * 2.09
    add_rect(sl, x, 4.28, 1.82, 0.88, fill=col)
    add_text(sl, step, x, 4.28, 1.82, 0.88, size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    if i < len(flow)-1:
        add_text(sl, "→", x+1.82, 4.48, 0.28, 0.42, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Bottom tagline
add_rect(sl, 0.4, 5.42, 12.53, 0.88, fill=RGBColor(0x0A, 0x2A, 0x50))
add_rect(sl, 0.4, 5.42, 0.1, 0.88, fill=CYAN)
add_text(sl, '"Given our real constraints — what is the best possible delivery schedule, '
             'and what does the unavoidable cost look like?"',
         0.65, 5.52, 12.1, 0.7, size=13, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "Slide 11 of 11", 0.3, 7.1, 3, 0.3, size=9, color=DGRAY)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "/home/user/projects/Order_Planned_System/OrderBasedPlanning_Presentation.pptx"
prs.save(out)
print("Saved:", out)
